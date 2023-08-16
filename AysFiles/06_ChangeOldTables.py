from pptx import Presentation

# PowerPoint dosyasını yükle
presentation = Presentation('Cam_details_template_(Draft_V01).pptx')

# Belirli bir hücrenin satır ve sütun indekslerini belirle
target_row = 2
target_column = 2
target_cell_value = "Pump type"

# Her slayttaki şekilleri kontrol et
for slide in presentation.slides:
    for shape in slide.shapes:
        if shape.has_table:
            
            table = shape.table
            # Hedef hücrenin içeriğini kontrol et
            if table.cell(0, 0).text == target_cell_value:
                print("Dostum hücreyi buldum yayy!! Değeri şuymuş: ", table.cell(0, 0).text)
                if target_row < len(table.rows) and target_column < len(table.columns):
                    cell = table.cell(target_row, target_column)
                    print("Önceki Değer: ", cell.text)
                    
                    # Hücrenin değerini güncelle
                    new_value = "Yeni Değer"
                    cell.text = new_value
                    print("Yeni Değer: ", cell.text)

# Değişiklikleri kaydet
presentation.save('06_ChangeOldTables.pptx')
