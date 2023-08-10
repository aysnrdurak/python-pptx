from pptx import Presentation
from pptx.chart.data import ChartData
import openpyxl

presentation = Presentation('Deneme.pptx')
wb_chart_details = openpyxl.Workbook()
ws_chart_details = wb_chart_details.active
ws_chart_details.title = "Chart Details"
ws_chart_details.append(["Slide ID", "Chart Type", "Chart Title"])

wb_chart_change = openpyxl.load_workbook("deneme.xlsx")
sheet = wb_chart_change['Tabelle1']
new_categories = [cell.value for cell in sheet['A'][1:]]
new_values = [cell.value for cell in sheet['B'][1:]]
new_values2 = [cell.value for cell in sheet['C'][1:]]

for slide in presentation.slides:
    for shape in slide.shapes:
        print(shape.shape_id)
        if shape.has_chart:
            chart = shape.chart

            #  Detect Charts
            chart_type = chart.chart_type
            chart_title = chart.chart_title.text_frame.text
            ws_chart_details.append([slide.slide_id, chart_type, chart_title])

            # Change Charts
            chart_data = ChartData()
            chart_data.categories = new_categories
            chart_data.add_series('Updated Data', new_values)
            chart_data.add_series('Updated Data', new_values2)
            chart.replace_data(chart_data)

wb_chart_details.save('deneme_chart_details.xlsx')
presentation.save('07_ChangeOldCharts.pptx')
