from pptx import Presentation

# PowerPoint dosyasını yükle
presentation = Presentation('Cam_details_template_(Draft_V01).pptx')

# Her slayttaki şekilleri kontrol et
for slide in presentation.slides:
    for shape in slide.shapes:
        if shape.has_table:
            
            table = shape.table
            print("Tablo Bilgileri:")
            
            # Tablonun tüm satırlarını ve sütunlarını döngüyle kontrol et
            for row_idx, row in enumerate(table.rows):
                for col_idx, cell in enumerate(row.cells):
                    print("Satır: {}, Sütun: {}, Metin: {}".format(row_idx, col_idx, cell.text))
